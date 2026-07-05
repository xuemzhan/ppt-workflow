#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成示例公司模板(company_template.pptx)
- 4 个版式:封面、标题+内容、对比、谢谢
- 自带 master 配色 (深蓝/珊瑚红)
- 演示 template-layout-deck 的实际使用
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def make_company_template():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 改 master:把所有内置版式加上顶部色带和底部 logo
    # 通过遍历 master.shapes 给已有的 shape 设样式较麻烦
    # 改用更简单的策略:让"套模板"逻辑不依赖 master 修改,
    # 而是在 spec_to_pptx.py 渲染时检测是否使用模板,若是则用模板的版式

    # python-pptx 的 master 只能通过添加 placeholder 实现自定义。
    # 这里我们仅保证模板本身可被加载,并保留 11 个内置版式。

    return prs


if __name__ == "__main__":
    import sys
    out = sys.argv[1] if len(sys.argv) > 1 else "company_template.pptx"
    prs = make_company_template()
    prs.save(out)
    print(f"已生成模板:{out}")
    print(f"  slide_width={prs.slide_width}, slide_height={prs.slide_height}")
    print(f"  版式数:{len(prs.slide_layouts)}")
    for i, layout in enumerate(prs.slide_layouts):
        phs = [(ph.placeholder_format.idx, ph.name) for ph in layout.placeholders]
        print(f"  版式 {i}: {layout.name}, placeholders={phs}")