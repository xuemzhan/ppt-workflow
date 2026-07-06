#!/usr/bin/env python3
"""PPT 工作流编排器 (PPT Workflow Orchestrator)

按文章《从"手搓 PPT"到"演示生产线":10 个 PPT Skills 构成的工作流》
定义的 10 步流水线,把资料 → outline → storyline → slides → spec →
draft → themed → final → 讲稿 → QA 报告 自动串起来。

特性:
- 错误传递:任何一步失败立即终止,不连环崩溃
- 幂等性:产物已存在自动跳过(类似 make),可用 --force 全量重跑
- 智能回退:对比栏/流程图数据不足时,自动回退到 bullet
- 真实讲稿:从 outline.md 自动生成草稿讲稿(非占位符)

用法:
    python orchestrate.py --input <input_dir> --output <output_dir> [--options]

input_dir 应包含:
    - theme.txt          演示主题(可选,默认推断)
    - audience.txt       受众(可选)
    - scenario.txt       场景(可选)
    - slides.json        演示骨架(若提供则跳过 outline-builder)
    - spec.json          完整 spec(若提供则跳过 markdown-to-slides)
    - template.pptx      公司模板(可选)
    - notes.md           讲稿(可选)
    - chart_data/        图表数据(可选)

output_dir 会生成:
    - outline.md
    - storyline.md
    - slides.md
    - spec.json
    - deck_draft.pptx
    - deck_themed.pptx       (若有模板)
    - deck_final.pptx
    - deck_final.pdf         (若有 LibreOffice)
    - speaker_notes.md
    - qa_report.md
    - charts/                (若有 chart_data)
    - workflow_log.md
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import argparse
import shutil
import subprocess
from datetime import datetime

from skills.shared.utils import (
    load_json,
    write_json,
)

SKILL_ROOT = Path(__file__).resolve().parent.parent / "skills"


def run_skill(skill_dir, script_name, *args):
    """调用一个 Skill 下的脚本(失败返回 False,继续执行)"""
    return _run_skill_impl(skill_dir, script_name, *args, die_on_error=False)


def run_skill_or_die(skill_dir, script_name, *args):
    """调用 Skill,失败立即终止整个流水线"""
    ok = _run_skill_impl(skill_dir, script_name, *args, die_on_error=True)
    return ok


def _run_skill_impl(skill_dir, script_name, *args, die_on_error):
    script_path = SKILL_ROOT / skill_dir / "scripts" / script_name
    if not script_path.exists():
        msg = f"  ⚠️  脚本不存在:{script_path}"
        print(msg)
        if die_on_error:
            sys.exit(1)
        return False
    cmd = [sys.executable, str(script_path), *args]
    print(f"  → {skill_dir}/{script_name} {' '.join(args)}")
    try:
        result = subprocess.run(cmd, check=False, capture_output=True, text=True)
        if result.returncode != 0:
            print(f"  ❌ 退出码 {result.returncode}")
            if result.stderr:
                # 只显示最后 500 字符,避免噪音
                stderr = result.stderr.strip()
                if len(stderr) > 500:
                    stderr = "..." + stderr[-500:]
                print(f"     stderr: {stderr}")
            if die_on_error:
                sys.exit(f"❌ 步骤失败:{skill_dir}/{script_name}")
            return False
        if result.stdout:
            for line in result.stdout.splitlines():
                if line.strip():
                    print(f"     | {line}")
        return True
    except Exception as e:
        print(f"  ❌ 调用失败:{e}")
        if die_on_error:
            sys.exit(1)
        return False


def read_or_default(path, default=""):
    if Path(path).exists():
        return Path(path).read_text(encoding="utf-8-sig").strip()
    return default


def list_skills():
    """列出所有可用 Skill 及其描述"""
    import re

    print("=" * 60)
    print("  可用 Skill 列表")
    print("=" * 60)
    skills_dir = SKILL_ROOT
    for skill_dir in sorted(skills_dir.iterdir()):
        if not skill_dir.is_dir():
            continue
        skill_md = skill_dir / "SKILL.md"
        if not skill_md.exists():
            continue
        skill_text = skill_md.read_text(encoding="utf-8-sig")
        # 提取 frontmatter
        name_match = re.search(r"^name:\s*(.+)$", skill_text, re.MULTILINE)
        desc_match = re.search(r"^description:\s*(.+)$", skill_text, re.MULTILINE)
        name = name_match.group(1) if name_match else skill_dir.name
        desc = desc_match.group(1) if desc_match else "(无描述)"
        # 提取脚本列表
        scripts_dir = skill_dir / "scripts"
        scripts = []
        if scripts_dir.exists():
            for s in sorted(scripts_dir.glob("*.py")):
                scripts.append(s.name)
        print(f"\n[{name}]")
        print(f"  描述:{desc[:100]}{'...' if len(desc) > 100 else ''}")
        if scripts:
            print(f"  脚本:{', '.join(scripts)}")
    print("\n" + "=" * 60)
    print("  用法:python orchestrate.py --input <dir> --output <dir>")
    print("  列出 Skill:python orchestrate.py --list")
    print("=" * 60)


def make_default_slides_json(theme, audience, scenario):
    """若用户没提供 slides.json,生成一个最小骨架"""
    return {
        "theme": theme,
        "audience": audience,
        "scenario": scenario,
        "slides": [
            {
                "section": "context",
                "title": f"为什么 {theme} 此刻重要",
                "points": ["背景 1", "背景 2"],
                "evidence": "数据",
                "chart_needed": False,
                "viz_type": "无",
                "source": "输入资料",
                "minutes": 1.5,
            },
            {
                "section": "challenge",
                "title": "我们面临的关键挑战",
                "points": ["挑战 1", "挑战 2"],
                "evidence": "案例",
                "chart_needed": False,
                "viz_type": "无",
                "source": "输入资料",
                "minutes": 1.5,
            },
            {
                "section": "insight",
                "title": "三个核心洞察",
                "points": ["洞察 1:...", "洞察 2:...", "洞察 3:..."],
                "evidence": "数据",
                "chart_needed": True,
                "chart_type": "柱状图",
                "viz_type": "无",
                "source": "输入资料",
                "minutes": 2,
            },
            {
                "section": "solution",
                "title": "我们的方案与路径",
                "points": ["方案 1", "方案 2", "方案 3"],
                "evidence": "案例",
                "chart_needed": False,
                "viz_type": "流程图",
                "source": "输入资料",
                "minutes": 2,
            },
            {
                "section": "evidence",
                "title": "数据与案例证明方案有效",
                "points": ["数据 1", "数据 2"],
                "evidence": "数据",
                "chart_needed": True,
                "chart_type": "折线图",
                "viz_type": "对比栏",
                "source": "输入资料",
                "minutes": 2,
            },
            {
                "section": "action",
                "title": "下一步:谁、什么时候、做什么",
                "points": ["责任 1 + 截止时间", "责任 2 + 截止时间"],
                "evidence": "行动",
                "chart_needed": False,
                "viz_type": "无",
                "source": "输入资料",
                "minutes": 1.5,
            },
        ],
    }


def should_skip(out_path, force):
    """幂等性:产物已存在且未指定 --force 时跳过"""
    if force:
        return False
    if out_path.exists() and out_path.stat().st_size > 0:  # noqa: SIM103
        return True
    return False


def convert_outline_to_storyline(outline_path):
    """从 outline.md 抽取每页标题/要点,生成草稿讲稿"""
    import re

    text = outline_path.read_text(encoding="utf-8-sig")
    # 提取每页的标题与"核心要点"下的子项(以 2 空格缩进的 - 开头)
    page_pattern = re.compile(
        r"###\s+页\s+(\d+):\s+(.+?)\n(.*?)(?=\n###\s+页\s+\d+:|\n##\s+[^#]|$)",
        re.DOTALL,
    )
    raw_pages = page_pattern.findall(text)
    pages = []
    for pn, title, body in raw_pages:
        # 在 body 中找 "核心要点:" 后面的 2 空格缩进 bullet
        pts_match = re.search(r"-\s*核心要点:\s*\n((?:\s{2,}-\s+.+\n?)+)", body)
        if pts_match:  # noqa: SIM108
            pts = re.findall(r"\s{2,}-\s+(.+)", pts_match.group(1))
        else:
            pts = []
        pages.append((pn, title.strip(), pts))
    audience_match = re.search(r"受众:(\S+)", text)
    audience_line = (
        f"> 受众:{audience_match.group(1)}\n" if audience_match else "> 受众:通用\n"
    )
    notes = [
        "# 演讲稿\n",
        audience_line,
        "> 自动生成:从 outline.md 派生(请人工润色)\n\n---\n",
    ]
    if not pages:
        # 兜底:从 markdown 抽取所有 ### 标题
        for m in re.finditer(r"^###\s+(.+)$", text, re.MULTILINE):
            title = m.group(1).strip()
            notes.append(f"\n## {title}\n")
            notes.append("**承接**: <请补充>\n")
            notes.append("**展开**: <请补充>\n")
            notes.append("**落地**: <请补充>\n")
            notes.append("**时长**: 1.5 分钟\n")
        return "\n".join(notes)

    for page_num, title, points in pages:
        title = title.strip()
        # points 已经是 list(由 page_pattern 解析时提取)
        notes.append(f"\n## 页 {page_num}: {title}\n")
        # 自动生成承接
        if page_num == "1":
            notes.append(f'**承接**: 大家好,今天汇报的主题是"{title}"。\n')
        else:
            prev_title = pages[int(page_num) - 2][1] if int(page_num) > 1 else ""
            notes.append(f'**承接**: 接上一页"{prev_title[:20]}",继续看...\n')
        # 自动展开(基于要点)
        if points:
            points_str = ";".join(points[:3])
            notes.append(f"**展开**: 重点讲三点:{points_str}。\n")
        else:
            notes.append("**展开**: <请补充>\n")
        # 自动落地
        if page_num == str(len(pages)):
            notes.append("**落地**: 总结以上,期待各位反馈。\n")
        else:
            notes.append("**落地**: <本节小结,引出下一页>\n")
        notes.append("**时长**: 1.5 分钟\n")
    return "\n".join(notes)


def orchestrate(
    input_dir,
    output_dir,
    theme=None,
    audience=None,
    scenario=None,
    template=None,
    force=False,
    skip_outline=False,
    skip_storyline=False,
    skip_template=False,
    skip_notes=False,
    skip_qa=False,
):
    input_dir = Path(input_dir).resolve()
    output_dir = Path(output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    if not theme:
        theme = read_or_default(input_dir / "theme.txt", "未命名演示")
    if not audience:
        audience = read_or_default(input_dir / "audience.txt", "通用受众")
    if not scenario:
        scenario = read_or_default(input_dir / "scenario.txt", "通用场景")

    print("=" * 60)
    print("  PPT 工作流编排器")
    print(f"  主题:{theme}")
    print(f"  受众:{audience}")
    print(f"  场景:{scenario}")
    print(f"  输入:{input_dir}")
    print(f"  输出:{output_dir}")
    print(f"  模式:{'全量重跑(--force)' if force else '增量(产物存在则跳过)'}")
    print(f"  时间:{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print("=" * 60)

    log_lines = [
        f"# PPT 工作流运行日志 - {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        f"主题:{theme} | 受众:{audience} | 场景:{scenario}",
        "",
    ]

    # Step 1: outline
    outline_path = output_dir / "outline.md"
    if should_skip(outline_path, force) and not skip_outline:
        print(f"\n[1/10] outline-builder — 跳过(已存在:{outline_path})")
        log_lines.append("⏭ Step 1 outline 跳过(产物存在)")
    else:
        print("\n[1/10] outline-builder — 生成 outline.md")
        slides_json = input_dir / "slides.json"
        if slides_json.exists():
            shutil.copy(slides_json, output_dir / "slides.json")
            data = load_json(slides_json)
        else:
            data = make_default_slides_json(theme, audience, scenario)
            write_json(output_dir / "slides.json", data)
        run_skill_or_die(
            "ppt-outline-builder",
            "build_outline.py",
            str(output_dir / "slides.json"),
            str(outline_path),
        )
        log_lines.append(f"✅ Step 1 outline → {outline_path}")

    # Step 2: storyline
    storyline_path = output_dir / "storyline.md"
    if should_skip(storyline_path, force) and not skip_storyline:
        print("\n[2/10] deck-storyline-designer — 跳过(已存在)")
        log_lines.append("⏭ Step 2 storyline 跳过")
    else:
        print("\n[2/10] deck-storyline-designer — 生成 storyline.md")
        run_skill_or_die(
            "deck-storyline-designer",
            "build_storyline.py",
            str(outline_path),
            "1",
            str(storyline_path),
        )
        log_lines.append(f"✅ Step 2 storyline → {storyline_path}")

    # Step 3: markdown-to-slides
    slides_md = output_dir / "slides.md"
    print("\n[3/10] markdown-to-slides — 生成 slides.md(由 storyline 简化)")
    if storyline_path.exists() and (force or not slides_md.exists()):
        text = storyline_path.read_text(encoding="utf-8-sig")
        import re

        pages = re.findall(r"###\s+页\s+(\d+):\s+(.+)", text)
        simple_lines = [f"# {theme}"]
        for i, (n, title) in enumerate(pages):
            if i > 0:
                simple_lines.append("\n---\n")
            simple_lines.append(f"\n## {title}\n")
        slides_md.write_text("\n".join(simple_lines), encoding="utf-8")
    print(f"     → {slides_md}")
    log_lines.append(f"✅ Step 3 slides.md → {slides_md}")

    # Step 4: spec.json
    spec_path = output_dir / "spec.json"
    print("\n[4/10] pptx-generator — 生成 spec.json")
    if force or not spec_path.exists():  # noqa: SIM102
        if outline_path.exists():
            import re

            text = outline_path.read_text(encoding="utf-8-sig")
            slides_data = load_json(output_dir / "slides.json")
            slides = [
                {
                    "type": "title",
                    "title": theme,
                    "subtitle": f"{audience} · {scenario}",
                }
            ]
            for s in slides_data.get("slides", []):
                title = s.get("title", "")
                points = s.get("points", [])
                viz_type = s.get("viz_type", "无")
                # P0-3 修复:智能回退,避免硬编码占位
                if viz_type == "对比栏" and len(points) >= 2:
                    mid = len(points) // 2 + 1
                    slides.append(
                        {
                            "type": "two_column",
                            "title": title,
                            "left": {"heading": "升级前", "items": points[:mid]},
                            "right": {"heading": "升级后", "items": points[mid:]},
                        }
                    )
                elif viz_type == "流程图" and points:
                    slides.append(
                        {
                            "type": "process",
                            "title": title,
                            "steps": [
                                {"title": p[:24], "description": p[24:48]}
                                for p in points[:6]
                            ],
                        }
                    )
                else:
                    slides.append(
                        {
                            "type": "bullet",
                            "title": title,
                            "bullets": [{"text": p} for p in points],
                        }
                    )
            slides.append(
                {"type": "thank_you", "title": "谢谢", "subtitle": "欢迎讨论"}
            )
            spec = {
                "theme": {
                    "primary_color": "1E2761",
                    "secondary_color": "CADCFC",
                    "accent_color": "F96167",
                    "header_font": "Microsoft YaHei",
                    "body_font": "Microsoft YaHei",
                },
                "slides": slides,
            }
            write_json(spec_path, spec)
            print(f"     → {spec_path}({len(slides)} 页)")
    log_lines.append(f"✅ Step 4 spec → {spec_path}")

    # Step 4.5: 验证 spec.json(失败立即终止)
    if spec_path.exists():
        print("\n[4.5/10] spec-validator — 验证 spec.json 结构")
        spec_valid_script = (
            SKILL_ROOT / "pptx-generator" / "scripts" / "validate_spec.py"
        )
        if spec_valid_script.exists():
            spec_check = subprocess.run(
                [sys.executable, str(spec_valid_script), str(spec_path)],
                capture_output=True,
                text=True,
            )
            if spec_check.returncode != 0:
                print("  ❌ spec.json 验证失败,终止流水线")
                if spec_check.stdout:
                    print(spec_check.stdout)
                if spec_check.stderr:
                    print(spec_check.stderr)
                sys.exit(1)
            else:
                for line in spec_check.stdout.splitlines():
                    if line.strip():
                        print(f"     | {line}")

    # Step 5: pptx-generator
    draft_path = output_dir / "deck_draft.pptx"
    if should_skip(draft_path, force):
        print("\n[5/10] pptx-generator — 跳过(已存在)")
        log_lines.append("⏭ Step 5 deck_draft 跳过")
    else:
        print("\n[5/10] pptx-generator — 生成 deck_draft.pptx")
        run_skill_or_die(
            "pptx-generator", "spec_to_pptx.py", str(spec_path), str(draft_path)
        )
        log_lines.append(f"✅ Step 5 deck_draft → {draft_path}")

    # Step 6: template-layout-deck
    themed_path = output_dir / "deck_themed.pptx"
    if template and Path(template).exists() and not skip_template:
        if should_skip(themed_path, force):
            print("\n[6/10] template-layout-deck — 跳过(已存在)")
            log_lines.append("⏭ Step 6 deck_themed 跳过")
        else:
            print(f"\n[6/10] template-layout-deck — 套模板:{template}")
            from copy import deepcopy

            from pptx import Presentation

            tmpl_prs = Presentation(str(template))
            draft_prs = Presentation(str(draft_path))
            for slide in draft_prs.slides:
                blank_layout = tmpl_prs.slide_layouts[6]
                new_slide = tmpl_prs.slides.add_slide(blank_layout)
                for shape in slide.shapes:
                    el = deepcopy(shape.element)
                    new_slide.shapes._spTree.insert_element_before(el, "p:extLst")
            tmpl_prs.save(str(themed_path))
            print(f"     → {themed_path}")
            log_lines.append(f"✅ Step 6 deck_themed → {themed_path}")
    else:
        print("\n[6/10] template-layout-deck — 跳过(无模板)")
        themed_path = draft_path

    # Step 7: deck_final
    final_path = output_dir / "deck_final.pptx"
    if should_skip(final_path, force):
        print("\n[7/10] deck_final.pptx — 跳过(已存在)")
        log_lines.append("⏭ Step 7 deck_final 跳过")
    else:
        shutil.copy(str(themed_path), str(final_path))
        print(f"\n[7/10] deck_final.pptx 已就绪 → {final_path}")
        log_lines.append(f"✅ Step 7 deck_final → {final_path}")

    # Step 8: speaker-notes(P1-4:从 outline 自动派生草稿)
    notes_path = output_dir / "speaker_notes.md"
    if not skip_notes:
        if force or not notes_path.exists():
            print("\n[8/10] speaker-notes-writer — 从 outline 生成草稿讲稿")
            notes_content = convert_outline_to_storyline(outline_path)
            notes_path.write_text(notes_content, encoding="utf-8")
            log_lines.append(f"✅ Step 8 notes → {notes_path}(从 outline 派生)")
        # 写入 PPTX 备注
        run_skill_or_die(
            "speaker-notes-writer",
            "notes_to_pptx.py",
            str(final_path),
            str(notes_path),
            str(final_path),
        )

    # Step 9: chart-slide-maker
    chart_dir = input_dir / "chart_data"
    if chart_dir.exists():
        out_chart_dir = output_dir / "charts"
        out_chart_dir.mkdir(exist_ok=True)
        for spec_file in chart_dir.glob("*.json"):
            png_path = out_chart_dir / (spec_file.stem + ".png")
            if should_skip(png_path, force):
                print(f"\n[9/10] chart-slide-maker — 跳过(已存在:{png_path.name})")
            else:
                print(f"\n[9/10] chart-slide-maker — {spec_file.name}")
                run_skill(
                    "chart-slide-maker",
                    "render_chart.py",
                    str(spec_file),
                    str(png_path),
                )
        log_lines.append(f"✅ Step 9 charts → {out_chart_dir}")

    # Step 10: deck-review-publisher(P1-5:扩展 QA)
    qa_path = output_dir / "qa_report.md"
    if not skip_qa:
        print("\n[10/10] deck-review-publisher — QA 质检")
        run_skill_or_die(
            "deck-review-publisher", "qa_check.py", str(final_path), str(qa_path)
        )
        # 尝试转 PDF(失败不阻塞)
        pdf_path = output_dir / "deck_final.pdf"
        if should_skip(pdf_path, force):
            print("  ⏭ PDF 已存在,跳过")
        else:
            try:
                _run_skill_impl(
                    "deck-review-publisher",
                    "soffice_convert.py",
                    str(final_path),
                    str(pdf_path),
                    die_on_error=False,
                )
                if pdf_path.exists():
                    log_lines.append(f"✅ Step 10 PDF → {pdf_path}")
            except Exception:
                print("  ⚠️  PDF 转换跳过(无 LibreOffice)")

        # 缩略图网格(零 LibreOffice 依赖,纯 python-pptx + Pillow)
        thumb_path = output_dir / "thumbnails.jpg"
        if not should_skip(thumb_path, force):
            try:
                _run_skill_impl(
                    "deck-review-publisher",
                    "render_thumbnails.py",
                    str(final_path),
                    str(thumb_path),
                    "--qa-report",
                    str(qa_path),
                    die_on_error=False,
                )
                if thumb_path.exists():
                    log_lines.append(f"✅ Step 10 thumbnails → {thumb_path}")
            except Exception as e:
                print(f"  ⚠️  缩略图生成失败:{e}")
        log_lines.append(f"✅ Step 10 qa → {qa_path}")

    log_path = output_dir / "workflow_log.md"
    log_path.write_text("\n".join(log_lines), encoding="utf-8")
    print("\n" + "=" * 60)
    print(f"  ✅ 全部完成!产出在:{output_dir}")
    print(f"  - {final_path.name}(终版 PPT)")
    if qa_path.exists():
        print(f"  - {qa_path.name}(质检报告)")
    if notes_path.exists():
        print(f"  - {notes_path.name}(演讲稿)")
    print("  - workflow_log.md(运行日志)")
    print("=" * 60)


def main():
    parser = argparse.ArgumentParser(description="PPT 工作流编排器")
    parser.add_argument("--input", help="输入目录")
    parser.add_argument("--output", help="输出目录")
    parser.add_argument("--theme", help="演示主题")
    parser.add_argument("--audience", help="受众")
    parser.add_argument("--scenario", help="场景")
    parser.add_argument("--template", help="公司/客户模板 .pptx 路径")
    parser.add_argument(
        "--force", action="store_true", help="强制全量重跑(忽略已有产物)"
    )
    parser.add_argument("--skip-outline", action="store_true")
    parser.add_argument("--skip-storyline", action="store_true")
    parser.add_argument("--skip-template", action="store_true")
    parser.add_argument("--skip-notes", action="store_true")
    parser.add_argument("--skip-qa", action="store_true")
    parser.add_argument("--list", action="store_true", help="列出所有可用 Skill 后退出")
    args = parser.parse_args()

    if args.list:
        list_skills()
        return

    if not args.input or not args.output:
        parser.error("--input 和 --output 是必需的(或使用 --list 查看所有 Skill)")

    orchestrate(
        input_dir=args.input,
        output_dir=args.output,
        theme=args.theme,
        audience=args.audience,
        scenario=args.scenario,
        template=args.template,
        force=args.force,
        skip_outline=args.skip_outline,
        skip_storyline=args.skip_storyline,
        skip_template=args.skip_template,
        skip_notes=args.skip_notes,
        skip_qa=args.skip_qa,
    )


if __name__ == "__main__":
    main()
