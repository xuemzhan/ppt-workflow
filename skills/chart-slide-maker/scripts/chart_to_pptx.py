#!/usr/bin/env python3
"""把 chart PNG 嵌入 PPTX 的一页
用法: chart_to_pptx.py <chart.png> <title> <interpretation> <output.pptx>
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import logging
from typing import Any

from skills.shared.logging_config import setup_logging
from skills.shared.utils import hex_to_rgb as _hex_to_rgb_tuple

setup_logging()
logger = logging.getLogger(__name__)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def hex_to_rgb(h: str) -> RGBColor:
    """Convert hex string to RGBColor,using shared/utils.hex_to_rgb internally."""
    r, g, b = _hex_to_rgb_tuple(h)
    return RGBColor(r, g, b)


def set_run(
    run: Any,
    text: str,
    size: int,
    bold: bool = False,
    color: RGBColor | None = None,
    font: str | None = None,
) -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if font:
        run.font.name = font


def main() -> None:
    parser = argparse.ArgumentParser(description="把 chart PNG 嵌入 PPTX")
    parser.add_argument("chart", help="图表 PNG 文件路径")
    parser.add_argument("title", help="幻灯片标题")
    parser.add_argument("interpretation", help="图表解读文字")
    parser.add_argument("output", help="输出 PPTX 文件路径")
    args = parser.parse_args()
    png_path = Path(args.chart)
    title = args.title
    interp = args.interpretation
    out_path = Path(args.output)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(
        p.add_run(),
        title,
        28,
        bold=True,
        color=hex_to_rgb("1E2761"),
        font="Microsoft YaHei",
    )

    slide.shapes.add_picture(
        str(png_path), Inches(0.5), Inches(1.3), height=Inches(5.0)
    )

    ix = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.8))
    itf = ix.text_frame
    ip = itf.paragraphs[0]
    set_run(
        ip.add_run(),
        "▎ " + interp,
        18,
        bold=True,
        color=hex_to_rgb("1E2761"),
        font="Microsoft YaHei",
    )

    prs.save(str(out_path))
    logger.info("已生成: %s", out_path)


if __name__ == "__main__":
    main()
