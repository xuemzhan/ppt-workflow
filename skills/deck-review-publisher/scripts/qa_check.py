#!/usr/bin/env python3
"""PPT QA 检查器:扫描 PPTX,输出问题清单

用法:
    python qa_check.py <deck.pptx> [output_report.md]

检查项:
- 占位符残留(lorem/ipsum/xxxx/click here 等)
- 硬编码占位词(对照 1/选项 A/示例文本/待填写 等)
- 标题中性(纯描述、无数字/方向)
- 重复文本(同一 shape 内 text 出现 ≥2 次)
- 字号过小(< 12pt 投影看不见)
- 备注完整性
- 页数合理性
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import logging
import re
from datetime import datetime

from skills.shared.logging_config import setup_logging

setup_logging()
logger = logging.getLogger(__name__)

from pptx import Presentation

PLACEHOLDER_PATTERNS = [
    r"lorem\s*ipsum",
    r"xxxx+",
    r"placeholder",
    r"click\s+here",
    r"点击此处",
    r"待补充",
    r"^\s*TODO\s*$",
    r"^\s*\[.*\]\s*$",
    r"^\s*<.*>\s*$",
    # 硬编码占位词(orchestrate.py v1 时代的产物)
    r"对照\s*\d+",
    r"选项\s*[A-Z]\b",
    r"示例文本",
    r"默认文本",
    r"sample\s*text",
    r"demo\s*data",
    # 中文占位
    r"待填写",
    r"此处输入",
    r"未填写",
    r"占位符",
    # LLM 常见占位
    r"\bxxx\b",
    r"\.\.\.\s*$",  # 行尾的省略号
]


NEUTRAL_TITLE_KEYWORDS = [
    "项目介绍",
    "项目概述",
    "项目背景",
    "项目目标",
    "工作汇报",
    "工作总结",
    "工作内容",
    "实验结果",
    "实验分析",
    "研究方法",
    "项目进度",
    "项目计划",
    "下一步",
    "团队介绍",
    "成员介绍",
    "产品介绍",
    "产品概述",
    "公司介绍",
    "公司简介",
]


def check_deck(pptx_path):
    prs = Presentation(str(pptx_path))
    issues = []
    total_slides = len(prs.slides)

    if total_slides < 5:
        issues.append(
            {
                "level": "🟡",
                "page": "-",
                "type": "page_count",
                "msg": "页数过少(" + str(total_slides) + "),可能未完成",
            }
        )
    elif total_slides > 30:
        issues.append(
            {
                "level": "🟡",
                "page": "-",
                "type": "page_count",
                "msg": "页数过多(" + str(total_slides) + "),可能需要拆分",
            }
        )

    for idx, slide in enumerate(prs.slides, 1):
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + "\n"

        for pat in PLACEHOLDER_PATTERNS:
            if re.search(pat, all_text, re.IGNORECASE | re.MULTILINE):
                issues.append(
                    {
                        "level": "🔴",
                        "page": idx,
                        "type": "placeholder",
                        "msg": "检测到占位符残留(模式:" + pat + ")",
                    }
                )

        first_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                first_text = shape.text_frame.text.strip()
                break

        if first_text and idx not in (1, total_slides):
            has_judgment = bool(
                re.search(r"\d", first_text)
                or re.search(r"[是为有将可达到]", first_text)
                or any(
                    kw in first_text
                    for kw in ["结论", "判断", "发现", "启示", "意味着"]
                )
            )
            if not has_judgment:
                for kw in NEUTRAL_TITLE_KEYWORDS:
                    if first_text == kw or first_text.startswith(kw):
                        issues.append(
                            {
                                "level": "🟡",
                                "page": idx,
                                "type": "neutral_title",
                                "msg": "标题偏中性:"
                                + first_text[:30]
                                + " - 建议改为带数字/方向的判断",
                            }
                        )
                        break

        has_note = False
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide


            notes = notes_slide.notes_text_frame.text.strip()
            has_note = bool(notes)
        if 2 <= idx <= total_slides - 1 and not has_note:
            issues.append(
                {
                    "level": "🟢",
                    "page": idx,
                    "type": "no_notes",
                    "msg": "本页缺少演讲者备注",
                }
            )

    # P1-5:重复文本检测(同一 shape 内 text 出现 ≥2 次)
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf_text = shape.text_frame.text.strip()
            if not tf_text or len(tf_text) < 10:
                continue
            lines = [raw_line.strip() for raw_line in tf_text.split("\n") if raw_line.strip()]
            if len(lines) >= 2:
                line_counts = {}
                for line in lines:
                    if len(line) >= 2:
                        line_counts[line] = line_counts.get(line, 0) + 1
                for line, count in line_counts.items():
                    if count >= 2:
                        issues.append(
                            {
                                "level": "🔴",
                                "page": idx,
                                "type": "duplicate_text",
                                "msg": "形状内重复文本(count="
                                + str(count)
                                + "):"
                                + line[:40],
                            }
                        )
                        break

    # P1-5:字号过小检测(< 12pt 投影看不见)
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is None:
                        continue
                    try:
                        size_pt = run.font.size.pt
                    except (AttributeError, TypeError):
                        continue
                    if size_pt > 0 and size_pt < 12:
                        text_preview = run.text[:20] if run.text else "(空)"
                        issues.append(
                            {
                                "level": "🟡",
                                "page": idx,
                                "type": "small_font",
                                "msg": "字号过小("
                                + str(round(size_pt, 1))
                                + "pt < 12pt):"
                                + text_preview
                                + " - 投影可能看不清",
                            }
                        )
                        break

    return issues


def main() -> None:
    parser = argparse.ArgumentParser(description="PPT QA 检查器")
    parser.add_argument("deck", help="PPTX 文件路径")
    parser.add_argument("output", nargs="?", help="输出报告 Markdown 文件路径")
    args = parser.parse_args()
    deck_path = Path(args.deck)
    out_path = Path(args.output) if args.output else None

    issues = check_deck(deck_path)
    prs = Presentation(str(deck_path))

    lines = [
        "# PPT 质检报告: " + deck_path.name,
        "",
        "> 生成时间:" + datetime.now().strftime("%Y-%m-%d %H:%M"),
        "> 总页数:" + str(len(prs.slides)),
        "> 问题总数:" + str(len(issues)),
        "",
        "## 问题统计",
        "",
        "- 🔴 阻塞:" + str(sum(1 for i in issues if i["level"] == "🔴")),
        "- 🟡 重要:" + str(sum(1 for i in issues if i["level"] == "🟡")),
        "- 🟢 优化:" + str(sum(1 for i in issues if i["level"] == "🟢")),
        "",
        "## 问题清单",
        "",
    ]

    if not issues:
        lines.append("✅ 未发现问题。")
    else:
        lines.append("| 级别 | 页 | 类型 | 描述 |")
        lines.append("|------|----|------|------|")
        for issue in issues:
            lines.append(
                "| "
                + issue["level"]
                + " | "
                + str(issue["page"])
                + " | "
                + issue["type"]
                + " | "
                + issue["msg"]
                + " |"
            )

    lines.extend(
        [
            "",
            "## 修复建议",
            "",
            "- 🔴 阻塞项:必须修复后才能发布",
            "- 🟡 重要项:尽量在演示前修复",
            "- 🟢 优化项:有时间就改",
            "",
            "## 修复路径",
            "",
            "| 类型 | 回退到哪个 Skill |",
            "|------|------------------|",
            "| placeholder / hardcoded_placeholder | pptx-generator / markdown-to-slides 重出 |",
            "| neutral_title | deck-storyline-designer 重写标题 |",
            "| no_notes | speaker-notes-writer 补讲稿 |",
            "| duplicate_text | 检查 spec.json,合并重复字段 |",
            "| small_font | pptx-generator 调大字号 |",
            "| 视觉问题 | visual-slide-designer 强化 / template-layout-deck 重套 |",
            "| page_count | ppt-outline-builder 重新规划 |",
            "",
        ]
    )

    out = "\n".join(lines)
    if out_path:
        out_path.write_text(out, encoding="utf-8")
        logger.info("已写入: %s", out_path)
    else:
        print(out)


if __name__ == "__main__":
    main()
