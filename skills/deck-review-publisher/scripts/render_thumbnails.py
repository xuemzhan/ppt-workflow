#!/usr/bin/env python3
"""PPT 缩略图网格生成器(零 LibreOffice 依赖)
读 PPTX,提取每页文本/颜色,生成 N x M 网格预览图。
不是精确像素还原,但能快速看出"哪页有问题"。

用法:
    python render_thumbnails.py <deck.pptx> <output.jpg> [--cols 4]
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import logging

from skills.shared.fonts import get_font_path
from skills.shared.logging_config import setup_logging

setup_logging()
logger = logging.getLogger(__name__)

logger = logging.getLogger(__name__)

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

THUMB_W, THUMB_H = 480, 270  # 16:9
LABEL_H = 30
COLS_DEFAULT = 4
BG = (250, 250, 250)
GRID_LINE = (220, 220, 220)
LABEL_BG = (30, 39, 97)
LABEL_FG = (255, 255, 255)


def get_font(size, bold=False):
    """Get a Pillow ImageFont using shared/fonts for cross-platform path discovery."""
    candidates_bold = ["msyhbd.ttc", "PingFang.ttc", "DejaVuSans-Bold.ttf"]
    candidates_reg = ["msyh.ttc", "PingFang.ttc", "DejaVuSans.ttf"]
    for name in candidates_bold if bold else candidates_reg:
        try:
            path = get_font_path(name)
            if path and path != "PIL_DEFAULT":
                return ImageFont.truetype(path, size)
        except Exception:
            continue
    try:
        path = get_font_path()
        if path and path != "PIL_DEFAULT":
            return ImageFont.truetype(path, size)
    except Exception:
        pass
    return ImageFont.load_default()


def emu_to_rgb(color_format):
    """提取主题色(简化:从 RGB 转换)"""
    try:
        return (color_format[0], color_format[1], color_format[2])
    except Exception:
        return (30, 39, 97)


def extract_bg_color(slide):
    """从 slide 背景提取主色(若有填充)"""
    try:
        fill = slide.background.fill
        if fill.type is not None and fill.fore_color.type is not None:
            try:
                rgb = fill.fore_color.rgb
                return (rgb[0], rgb[1], rgb[2])
            except Exception:
                pass
    except Exception:
        pass
    return (255, 255, 255)


def truncate(text, max_chars):
    text = text.strip()
    if len(text) > max_chars:
        return text[:max_chars] + "..."
    return text


def render_one_slide(slide):
    """把一页 PPT 渲染成 480x270 缩略图(简化版)"""
    img = Image.new("RGB", (THUMB_W, THUMB_H), extract_bg_color(slide))
    d = ImageDraw.Draw(img)

    # 取前几个文本
    texts = []
    title = ""
    bullets = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            t = "".join(r.text for r in p.runs).strip()
            if t:
                texts.append(t)
        if texts:
            break  # 第一个 text_frame 是标题/容器

    # 用 slide.shapes 顺序扫
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs).strip()
                if t:
                    texts.append(t)
    if not texts:
        texts = ["(空)"]

    title = truncate(texts[0], 30)
    bullets = [truncate(t, 35) for t in texts[1:5]]

    # 标题
    title_font = get_font(20, bold=True)
    d.text(
        (15, 15),
        title,
        fill=(255, 255, 255)
        if extract_bg_color(slide) != (255, 255, 255)
        else (30, 39, 97),
        font=title_font,
    )

    # bullet 列表
    bullet_font = get_font(13)
    y = 60
    for b in bullets:
        d.text(
            (20, y),
            "• " + b,
            fill=(60, 60, 60)
            if extract_bg_color(slide) == (255, 255, 255)
            else (240, 240, 240),
            font=bullet_font,
        )
        y += 22
        if y > THUMB_H - 30:
            break

    # 装饰条
    bar = d.rectangle([(0, 0), (4, THUMB_H)], fill=(249, 97, 103))
    return img


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("deck", help="输入 PPTX")
    parser.add_argument("output", help="输出 JPG/PNG")
    parser.add_argument("--cols", type=int, default=COLS_DEFAULT)
    parser.add_argument(
        "--qa-report", help="可选 qa_report.md 路径,有则在缩略图上加警告标记"
    )
    args = parser.parse_args()

    deck_path = Path(args.deck)
    out_path = Path(args.output)
    if not deck_path.exists():
        logger.error("%s 不存在", deck_path)
        sys.exit(1)

    prs = Presentation(str(deck_path))
    slides = list(prs.slides)
    n = len(slides)
    if n == 0:
        logger.error("PPTX 无幻灯片")
        sys.exit(1)

    cols = min(args.cols, n)
    rows = (n + cols - 1) // cols
    cell_w = THUMB_W
    cell_h = THUMB_H + LABEL_H
    grid = Image.new("RGB", (cols * cell_w, rows * cell_h), BG)
    d = ImageDraw.Draw(grid)
    label_font = get_font(18, bold=True)

    # 解析 qa_report.md: {page_num: highest_level}
    qa_issues = {}
    if args.qa_report and Path(args.qa_report).exists():
        import re

        for line in Path(args.qa_report).read_text(encoding="utf-8").splitlines():
            m = re.match(r"\| (\S+)\s+\|\s*(\d+)\s*\|", line)
            if m:
                level, page = m.group(1), int(m.group(2))
                if page == "-":
                    continue
                # 优先级:🔴 > 🟡 > 🟢
                rank = {"🔴": 3, "🟡": 2, "🟢": 1}.get(level, 0)
                if page not in qa_issues or qa_issues[page] < rank:
                    qa_issues[page] = rank
    if qa_issues:
        logger.warning("QA 警告: 共 %d 页有问题", len(qa_issues))
    warn_font = get_font(28, bold=True)

    for i, slide in enumerate(slides):
        row = i // cols
        col = i % cols
        thumb = render_one_slide(slide)
        x = col * cell_w
        y = row * cell_h
        grid.paste(thumb, (x, y))
        # 边线颜色:有 QA 问题时用红/黄
        outline_color = GRID_LINE
        if (i + 1) in qa_issues:
            rank = qa_issues[i + 1]
            if rank >= 3:
                outline_color = (200, 30, 30)
                warn_color = (220, 30, 30)
            elif rank >= 2:
                outline_color = (220, 180, 30)
                warn_color = (240, 200, 30)
            else:
                outline_color = (180, 180, 180)
                warn_color = (180, 180, 180)
            d.rectangle(
                [(x, y), (x + cell_w, y + cell_h)], outline=outline_color, width=4
            )
            # 右上角警告圈(三种等级都画,统一"!"图标)
            d.ellipse(
                [(x + cell_w - 50, y + 5), (x + cell_w - 5, y + 50)],
                fill=warn_color,
                outline=(80, 80, 80),
                width=2,
            )
            d.text((x + cell_w - 38, y + 10), "!", fill=(255, 255, 255), font=warn_font)
        else:
            d.rectangle([(x, y), (x + cell_w, y + cell_h)], outline=GRID_LINE, width=1)
        # 标签
        d.rectangle([(x, y + THUMB_H), (x + cell_w, y + cell_h)], fill=LABEL_BG)
        d.text((x + 10, y + THUMB_H + 5), f"页 {i + 1}", fill=LABEL_FG, font=label_font)

    grid.save(str(out_path), "JPEG", quality=85)
    logger.info("已生成缩略图网格: %s (%d 页, %dx%d)", out_path, n, cols, rows)


if __name__ == "__main__":
    main()
