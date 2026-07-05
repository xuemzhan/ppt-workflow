#!/usr/bin/env python3
"""Markdown → PPTX 转换器(零外部依赖,基于 python-pptx)

用法:
    python md_to_pptx.py <input.md> <output.pptx> [template_profile.json]

支持语法:
    # 主题
    --- 分页
    ## 标题
    - 要点
    > 引用
    `数字`  大数字
    **粗体**
    <!-- _class: cover|hook|stat|compare|timeline|quote|closing -->
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import logging
import re
from typing import Any

from skills.shared.logging_config import setup_logging
from skills.shared.utils import hex_to_rgb as _hex_to_rgb_tuple
from skills.shared.utils import load_json

setup_logging()
logger = logging.getLogger(__name__)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DEFAULT_THEME: dict[str, str] = {
    "primary_color": "1E2761",
    "secondary_color": "CADCFC",
    "accent_color": "F96167",
    "text_dark": "212121",
    "text_light": "FFFFFF",
    "header_font": "Microsoft YaHei",
    "body_font": "Microsoft YaHei",
}


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert hex color string to RGBColor,using shared/utils.hex_to_rgb internally."""
    r, g, b = _hex_to_rgb_tuple(hex_str)
    return RGBColor(r, g, b)


def parse_md_slides(md_text: str) -> tuple[str, list[dict[str, Any]]]:
    theme_match = re.match(r"^#\s+(.+)$", md_text.split("\n", 1)[0] if md_text else "")
    theme = theme_match.group(1).strip() if theme_match else "未命名演示"
    body = re.sub(r"^#\s+.+\n", "", md_text, count=1)
    raw_pages = re.split(r"^---\s*$", body, flags=re.MULTILINE)
    slides = []
    for raw in raw_pages:
        raw = raw.strip()
        if not raw:
            continue
        slides.append(parse_one_slide(raw))
    return theme, slides


def parse_one_slide(raw: str) -> dict[str, Any]:
    slide: dict[str, Any] = {
        "title": "",
        "points": [],
        "quote": None,
        "stat": None,
        "class_hint": None,
    }
    for line in raw.split("\n"):
        s = line.rstrip()
        m = re.match(r"<!--\s*_class:\s*(\S+)\s*-->", s)
        if m:
            slide["class_hint"] = m.group(1)
            continue
        m = re.match(r"^##\s+(.+)$", s)
        if m:
            slide["title"] = m.group(1).strip()
            continue
        m = re.match(r"^[-*]\s+(.+)$", s)
        if m:
            slide["points"].append(m.group(1).strip())
            continue
        m = re.match(r"^>\s+(.+)$", s)
        if m:
            slide["quote"] = m.group(1).strip()
            continue
        m = re.match(r"^`(\d[\d,.]*)\`\s*(.*)$", s)
        if m:
            slide["stat"] = {"number": m.group(1), "label": m.group(2).strip()}
            continue
        if s and not slide["title"]:
            slide["title"] = s
        elif s:
            slide["points"].append(s)
    return slide


def detect_class(slide: dict[str, Any], idx: int, total: int) -> str:
    if slide.get("class_hint"):
        return slide["class_hint"]
    if idx == 0:
        return "cover"
    if idx == total - 1:
        return "closing"
    if slide.get("stat"):
        return "stat"
    if slide.get("quote") and not slide.get("points"):
        return "quote"
    if len(slide.get("points", [])) == 2 and any(
        kw in " ".join(slide["points"]).lower() for kw in ["vs", "对比"]
    ):
        return "compare"
    if any(re.search(r"Q[1-4]|阶段|\d{4}年", p) for p in slide.get("points", [])):
        return "timeline"
    return "content"


def set_run(
    run: Any,
    text: str,
    size: int,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor | None = None,
    font: str | None = None,
) -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    if font:
        run.font.name = font


def add_cover(
    prs: Presentation, theme_cfg: dict[str, str], slide: dict[str, Any]
) -> None:
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    fill = blank.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme_cfg["primary_color"])
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(
        p.add_run(),
        slide["title"],
        54,
        bold=True,
        color=hex_to_rgb(theme_cfg["text_light"]),
        font=theme_cfg["header_font"],
    )
    if slide.get("points"):
        sub = blank.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.33), Inches(1))
        sf = sub.text_frame
        sf.word_wrap = True
        sp = sf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        set_run(
            sp.add_run(),
            slide["points"][0],
            20,
            color=hex_to_rgb(theme_cfg["secondary_color"]),
            font=theme_cfg["body_font"],
        )


def add_content(
    prs: Presentation, theme_cfg: dict[str, str], slide: dict[str, Any]
) -> None:
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(
        p.add_run(),
        slide["title"],
        36,
        bold=True,
        color=hex_to_rgb(theme_cfg["primary_color"]),
        font=theme_cfg["header_font"],
    )

    bar = blank.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.55), Inches(0.6), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme_cfg["accent_color"])
    bar.line.fill.background()

    if slide["points"]:
        bx = blank.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(11.8), Inches(5))
        bf = bx.text_frame
        bf.word_wrap = True
        for i, point in enumerate(slide["points"]):
            p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(10)
            segments = re.split(r"(\*\*[^*]+\*\*)", point)
            for seg in segments:
                if seg.startswith("**") and seg.endswith("**"):
                    set_run(
                        p.add_run(),
                        seg[2:-2],
                        20,
                        bold=True,
                        color=hex_to_rgb(theme_cfg["accent_color"]),
                        font=theme_cfg["body_font"],
                    )
                elif seg:
                    set_run(
                        p.add_run(),
                        seg,
                        20,
                        color=hex_to_rgb(theme_cfg["text_dark"]),
                        font=theme_cfg["body_font"],
                    )


def add_stat(
    prs: Presentation, theme_cfg: dict[str, str], slide: dict[str, Any]
) -> None:
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    fill = blank.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme_cfg["primary_color"])
    if slide.get("stat"):
        tx = blank.shapes.add_textbox(
            Inches(0.5), Inches(2.0), Inches(12.33), Inches(3)
        )
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(
            p.add_run(),
            slide["stat"]["number"],
            96,
            bold=True,
            color=hex_to_rgb(theme_cfg["accent_color"]),
            font=theme_cfg["header_font"],
        )
        lx = blank.shapes.add_textbox(
            Inches(0.5), Inches(5.0), Inches(12.33), Inches(1.5)
        )
        lf = lx.text_frame
        lp = lf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        set_run(
            lp.add_run(),
            slide["stat"]["label"] or slide.get("title", ""),
            28,
            color=hex_to_rgb(theme_cfg["text_light"]),
            font=theme_cfg["body_font"],
        )


def add_quote(
    prs: Presentation, theme_cfg: dict[str, str], slide: dict[str, Any]
) -> None:
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    fill = blank.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme_cfg["secondary_color"])
    tx = blank.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(3))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(
        p.add_run(),
        '"' + (slide.get("quote") or slide.get("title", "")) + '"',
        40,
        italic=True,
        color=hex_to_rgb(theme_cfg["primary_color"]),
        font=theme_cfg["header_font"],
    )


def add_compare(
    prs: Presentation, theme_cfg: dict[str, str], slide: dict[str, Any]
) -> None:
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    set_run(
        p.add_run(),
        slide["title"],
        32,
        bold=True,
        color=hex_to_rgb(theme_cfg["primary_color"]),
        font=theme_cfg["header_font"],
    )

    if len(slide["points"]) >= 2:
        left_text, right_text = slide["points"][0], slide["points"][1]
    else:
        left_text, right_text = "", ""

    lb = blank.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.0), Inches(5.8), Inches(4.5)
    )
    lb.fill.solid()
    lb.fill.fore_color.rgb = hex_to_rgb(theme_cfg["primary_color"])
    lb.line.fill.background()
    ltx = lb.text_frame
    ltx.word_wrap = True
    lp = ltx.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    set_run(
        lp.add_run(),
        left_text,
        24,
        color=hex_to_rgb(theme_cfg["text_light"]),
        font=theme_cfg["body_font"],
    )

    rb = blank.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.83), Inches(2.0), Inches(5.8), Inches(4.5)
    )
    rb.fill.solid()
    rb.fill.fore_color.rgb = hex_to_rgb(theme_cfg["accent_color"])
    rb.line.fill.background()
    rtx = rb.text_frame
    rtx.word_wrap = True
    rp = rtx.paragraphs[0]
    rp.alignment = PP_ALIGN.CENTER
    set_run(
        rp.add_run(),
        right_text,
        24,
        color=hex_to_rgb(theme_cfg["text_light"]),
        font=theme_cfg["body_font"],
    )


def add_timeline(
    prs: Presentation, theme_cfg: dict[str, str], slide: dict[str, Any]
) -> None:
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    set_run(
        p.add_run(),
        slide["title"],
        32,
        bold=True,
        color=hex_to_rgb(theme_cfg["primary_color"]),
        font=theme_cfg["header_font"],
    )

    line = blank.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.6), Inches(12.33), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(theme_cfg["primary_color"])
    line.line.fill.background()

    points = slide.get("points", [])
    n = max(len(points), 1)
    step = 12.33 / max(n, 1)
    for i, p_text in enumerate(points):
        cx = Inches(0.5 + step * (i + 0.5))
        dot = blank.shapes.add_shape(
            MSO_SHAPE.OVAL, cx - Inches(0.15), Inches(3.45), Inches(0.3), Inches(0.3)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb(theme_cfg["accent_color"])
        dot.line.fill.background()
        ltx = blank.shapes.add_textbox(
            cx - Inches(step / 2), Inches(2.0), Inches(step), Inches(1.3)
        )
        ltf = ltx.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        set_run(
            lp.add_run(),
            p_text[:30],
            14,
            bold=True,
            color=hex_to_rgb(theme_cfg["primary_color"]),
            font=theme_cfg["body_font"],
        )


def add_closing(
    prs: Presentation, theme_cfg: dict[str, str], slide: dict[str, Any]
) -> None:
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    fill = blank.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme_cfg["primary_color"])
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.33), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(
        p.add_run(),
        slide["title"] or "谢谢",
        60,
        bold=True,
        color=hex_to_rgb(theme_cfg["text_light"]),
        font=theme_cfg["header_font"],
    )
    if slide.get("points"):
        stx = blank.shapes.add_textbox(
            Inches(0.5), Inches(5.0), Inches(12.33), Inches(2)
        )
        stf = stx.text_frame
        for i, pt_text in enumerate(slide["points"]):
            sp = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
            sp.alignment = PP_ALIGN.CENTER
            set_run(
                sp.add_run(),
                pt_text,
                18,
                color=hex_to_rgb(theme_cfg["secondary_color"]),
                font=theme_cfg["body_font"],
            )


RENDERERS = {
    "cover": add_cover,
    "content": add_content,
    "hook": add_content,
    "stat": add_stat,
    "quote": add_quote,
    "compare": add_compare,
    "timeline": add_timeline,
    "closing": add_closing,
}


def main() -> None:
    parser = argparse.ArgumentParser(description="Markdown 转 PPTX")
    parser.add_argument("input", help="输入 Markdown 文件路径")
    parser.add_argument("output", help="输出 PPTX 文件路径")
    parser.add_argument("profile", nargs="?", help="模板配置 JSON 文件路径")
    args = parser.parse_args()
    md_path = Path(args.input)
    out_path = Path(args.output)
    theme_cfg = dict(DEFAULT_THEME)
    if args.profile:
        prof = load_json(Path(args.profile))
        theme_cfg.update(prof)

    md_text = md_path.read_text(encoding="utf-8-sig")
    _, slides = parse_md_slides(md_text)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, slide in enumerate(slides):
        cls = detect_class(slide, idx, len(slides))
        renderer = RENDERERS.get(cls, add_content)
        renderer(prs, theme_cfg, slide)

    prs.save(str(out_path))
    logger.info("已生成: %s (%d 页)", out_path, len(slides))


if __name__ == "__main__":
    main()
