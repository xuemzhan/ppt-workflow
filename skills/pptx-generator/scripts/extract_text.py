#!/usr/bin/env python3
"""提取 PPTX 全部文本(标题 + 正文 + 备注)
用法: extract_text.py <file.pptx>
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import logging
import sys
from pathlib import Path

from pptx import Presentation
from skills.shared.logging_config import setup_logging

setup_logging()
logger = logging.getLogger(__name__)


def main() -> None:
    parser = argparse.ArgumentParser(description="提取 PPTX 全部文本")
    parser.add_argument("input", help="PPTX 文件路径")
    args = parser.parse_args()

    prs = Presentation(args.input)
    for idx, slide in enumerate(prs.slides, 1):
        print(f"\n===== 页 {idx} =====")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs)
                    if t.strip():
                        print(t)
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes = notes_slide.notes_text_frame.text
            if notes.strip():
                print(f"  [备注] {notes}")


if __name__ == "__main__":
    main()
