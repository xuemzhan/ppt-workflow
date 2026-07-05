#!/usr/bin/env python3
"""合并多个 PPTX
用法: merge_pptx.py <out.pptx> <in1.pptx> <in2.pptx> [...]
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import logging
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from skills.shared.logging_config import setup_logging

setup_logging()
logger = logging.getLogger(__name__)


def main() -> None:
    parser = argparse.ArgumentParser(description="合并多个 PPTX 文件")
    parser.add_argument("output", help="输出 PPTX 文件路径")
    parser.add_argument("inputs", nargs="+", help="输入 PPTX 文件路径")
    args = parser.parse_args()

    out_path = Path(args.output)
    merged = Presentation(str(args.inputs[0]))

    for src in args.inputs[1:]:
        prs = Presentation(str(src))
        for slide in prs.slides:
            blank_layout = merged.slide_layouts[6]
            new_slide = merged.slides.add_slide(blank_layout)
            for shape in slide.shapes:
                el = deepcopy(shape.element)
                new_slide.shapes._spTree.insert_element_before(el, "p:extLst")

    merged.save(str(out_path))
    logger.info("已合并到 %s", out_path)


if __name__ == "__main__":
    main()
