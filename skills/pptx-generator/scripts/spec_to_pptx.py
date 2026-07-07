#!/usr/bin/env python3
"""PPTX 规格生成器:从 JSON spec 精确生成 PPTX

用法: spec_to_pptx.py <spec.json> <output.pptx>

支持的 slide type:
    title / section / bullet / two_column / three_column
    stat / quote / process / table / timeline / thank_you
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import logging
from typing import Any

from skills.shared.logging_config import setup_logging
from skills.shared.utils import hex_to_rgb as _hex_to_rgb_tuple
from skills.shared.utils import load_json

setup_logging()
logger = logging.getLogger(__name__)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(h: str) -> RGBColor:
    """Convert hex string to RGBColor,using shared/utils.hex_to_rgb internally."""
    r, g, b = _hex_to_rgb_tuple(h)
    return RGBColor(r, g, b)


def set_run(
    run: Any,
    text: str,
    size: int,
    bold: bool = False,
    color: RGBColor | None = None,
    font: str | None = None,
    italic: bool = False,
) -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    if font:
        run.font.name = font


def add_title(prs: Presentation, theme: dict[str, str], s: dict[str, Any]) -> None:
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    fill = blank.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme.get("primary_color", "1E2761"))
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(
        p.add_run(),
        s["title"],
        54,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        font=theme.get("header_font", "Microsoft YaHei"),
    )
    if s.get("subtitle"):
        stx = blank.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(12.33), Inches(1)
        )
        stf = stx.text_frame
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        set_run(
            sp.add_run(),
            s["subtitle"],
            24,
            color=RGBColor(0xCA, 0xDC, 0xFC),
            font=theme.get("body_font", "Microsoft YaHei"),
        )
    if s.get("author") or s.get("date"):
        atx = blank.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.6)
        )
        atf = atx.text_frame
        ap = atf.paragraphs[0]
        ap.alignment = PP_ALIGN.CENTER
        text = " · ".join(filter(None, [s.get("author"), s.get("date")]))
        set_run(
            ap.add_run(),
            text,
            14,
            color=RGBColor(0xCA, 0xDC, 0xFC),
            font=theme.get("body_font", "Microsoft YaHei"),
        )


def add_section(prs: Presentation, theme: dict[str, str], s: dict[str, Any]) -> None:
    # type: ignore[reportGeneralTypeIssues]
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    fill = blank.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme.get("primary_color", "1E2761"))
    if s.get("section_number"):
        ntx = blank.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.33), Inches(2)
        )
        ntf = ntx.text_frame
        np_ = ntf.paragraphs[0]
        np_.alignment = PP_ALIGN.CENTER
        set_run(
            np_.add_run(),
            s["section_number"],
            80,
            bold=True,
            color=hex_to_rgb(theme.get("accent_color", "F96167")),
            font=theme.get("header_font", "Microsoft YaHei"),
        )
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.33), Inches(1.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(
        p.add_run(),
        s["title"],
        36,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        font=theme.get("header_font", "Microsoft YaHei"),
    )


def add_bullet(prs: Presentation, theme: dict[str, str], s: dict[str, Any]) -> None:
    # type: ignore[reportGeneralTypeIssues]
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(
        p.add_run(),
        s["title"],
        36,
        bold=True,
        color=hex_to_rgb(theme.get("primary_color", "1E2761")),
        font=theme.get("header_font", "Microsoft YaHei"),
    )
    bar = blank.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.55), Inches(0.6), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme.get("accent_color", "F96167"))
    bar.line.fill.background()

    if s.get("bullets"):
        bx = blank.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(11.8), Inches(5))
        bf = bx.text_frame
        bf.word_wrap = True
        for i, b in enumerate(s["bullets"]):
            p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(12)
            indent = "    " * b.get("level", 0)
            set_run(
                p.add_run(),
                indent + "• " + b["text"],
                22,
                color=hex_to_rgb("212121"),
                font=theme.get("body_font", "Microsoft YaHei"),
            )

    if s.get("note"):
        blank.notes_slide.notes_text_frame.text = s["note"]


def add_two_column(prs, theme, s):
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    set_run(
        p.add_run(),
        s["title"],
        32,
        bold=True,
        color=hex_to_rgb(theme.get("primary_color", "1E2761")),
        font=theme.get("header_font", "Microsoft YaHei"),
    )
    for idx, key in enumerate(["left", "right"]):
        col = s.get(key, {})
        cx = Inches(0.7 + idx * 6.2)
        col_box = blank.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, cx, Inches(2.0), Inches(5.8), Inches(4.5)
        )
        col_box.fill.solid()
        col_box.fill.fore_color.rgb = hex_to_rgb(
            theme.get("primary_color", "1E2761")
            if idx == 0
            else theme.get("accent_color", "F96167")
        )
        col_box.line.fill.background()
        ctf = col_box.text_frame
        ctf.word_wrap = True
        ch = ctf.paragraphs[0]
        ch.alignment = PP_ALIGN.CENTER
        set_run(
            ch.add_run(),
            col.get("heading", ""),
            26,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            font=theme.get("header_font", "Microsoft YaHei"),
        )
        for item in col.get("items", []):
            ip = ctf.add_paragraph()
            ip.alignment = PP_ALIGN.LEFT
            ip.space_before = Pt(8)
            set_run(
                ip.add_run(),
                "• " + item,
                16,
                color=RGBColor(0xFF, 0xFF, 0xFF),
                font=theme.get("body_font", "Microsoft YaHei"),
            )


def add_stat(prs: Presentation, theme: dict[str, str], s: dict[str, Any]) -> None:
    # type: ignore[reportGeneralTypeIssues]
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(
        p.add_run(),
        s["title"],
        32,
        bold=True,
        color=hex_to_rgb(theme.get("primary_color", "1E2761")),
        font=theme.get("header_font", "Microsoft YaHei"),
    )
    stats = s.get("stats", [])
    if not stats:
        return
    width = 12.33 / len(stats)
    for i, st in enumerate(stats):
        cx = Inches(0.5 + i * width)
        box = blank.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, cx, Inches(2.0), Inches(width - 0.3), Inches(4.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = hex_to_rgb(theme.get("secondary_color", "CADCFC"))
        box.line.fill.background()
        btf = box.text_frame
        btf.word_wrap = True
        vp = btf.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        set_run(
            vp.add_run(),
            st["value"],
            60,
            bold=True,
            color=hex_to_rgb(theme.get("accent_color", "F96167")),
            font=theme.get("header_font", "Microsoft YaHei"),
        )
        lp = btf.add_paragraph()
        lp.alignment = PP_ALIGN.CENTER
        lp.space_before = Pt(12)
        set_run(
            lp.add_run(),
            st.get("label", ""),
            18,
            bold=True,
            color=hex_to_rgb(theme.get("primary_color", "1E2761")),
            font=theme.get("body_font", "Microsoft YaHei"),
        )
        dp = btf.add_paragraph()
        dp.alignment = PP_ALIGN.CENTER
        dp.space_before = Pt(8)
        set_run(
            dp.add_run(),
            st.get("description", ""),
            14,
            color=hex_to_rgb("212121"),
            font=theme.get("body_font", "Microsoft YaHei"),
        )


def add_quote(prs: Presentation, theme: dict[str, str], s: dict[str, Any]) -> None:
    # type: ignore[reportGeneralTypeIssues]
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    fill = blank.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme.get("secondary_color", "CADCFC"))
    tx = blank.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(3))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(
        p.add_run(),
        '"' + s["text"] + '"',
        40,
        italic=True,
        color=hex_to_rgb(theme.get("primary_color", "1E2761")),
        font=theme.get("header_font", "Microsoft YaHei"),
    )
    if s.get("attribution"):
        atx = blank.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1))
        atf = atx.text_frame
        ap = atf.paragraphs[0]
        ap.alignment = PP_ALIGN.CENTER
        set_run(
            ap.add_run(),
            s["attribution"],
            18,
            color=hex_to_rgb(theme.get("primary_color", "1E2761")),
            font=theme.get("body_font", "Microsoft YaHei"),
        )


def add_process(prs: Presentation, theme: dict[str, str], s: dict[str, Any]) -> None:
    # type: ignore[reportGeneralTypeIssues]
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    set_run(
        p.add_run(),
        s["title"],
        32,
        bold=True,
        color=hex_to_rgb(theme.get("primary_color", "1E2761")),
        font=theme.get("header_font", "Microsoft YaHei"),
    )
    steps = s.get("steps", [])
    if not steps:
        return
    n = len(steps)
    box_w = 12.33 / n
    for i, step in enumerate(steps):
        cx = Inches(0.5 + i * box_w + 0.1)
        circle = blank.shapes.add_shape(
            MSO_SHAPE.OVAL,
            cx + Inches(box_w / 2 - 0.4),
            Inches(2.0),
            Inches(0.8),
            Inches(0.8),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = hex_to_rgb(theme.get("accent_color", "F96167"))
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        set_run(
            cp.add_run(),
            str(i + 1),
            28,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            font=theme.get("header_font", "Microsoft YaHei"),
        )
        ttx = blank.shapes.add_textbox(
            cx, Inches(3.0), Inches(box_w - 0.2), Inches(0.8)
        )
        ttf = ttx.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        title_text = step.get("title", "")
        desc_text = step.get("description", "") or title_text
        set_run(
            tp.add_run(),
            title_text,
            16,
            bold=True,
            color=hex_to_rgb(theme.get("primary_color", "1E2761")),
            font=theme.get("header_font", "Microsoft YaHei"),
        )
        dtx = blank.shapes.add_textbox(
            cx, Inches(3.8), Inches(box_w - 0.2), Inches(2.0)
        )
        dtf = dtx.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.alignment = PP_ALIGN.CENTER
        set_run(
            dp.add_run(),
            desc_text,
            12,
            color=hex_to_rgb("212121"),
            font=theme.get("body_font", "Microsoft YaHei"),
        )
        if i < n - 1:
            arrow = blank.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                cx + Inches(box_w - 0.3),
                Inches(2.32),
                Inches(0.3),
                Inches(0.16),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = hex_to_rgb(theme.get("primary_color", "1E2761"))
            arrow.line.fill.background()


def add_table(prs, theme, s):
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    set_run(
        p.add_run(),
        s["title"],
        32,
        bold=True,
        color=hex_to_rgb(theme.get("primary_color", "1E2761")),
        font=theme.get("header_font", "Microsoft YaHei"),
    )
    headers = s.get("headers", [])
    rows = s.get("rows", [])
    n_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
    n_rows = len(rows) + (1 if headers else 0)
    if n_rows == 0:
        return
    table_shape = blank.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(2.0), Inches(12.33), Inches(4.5)
    )
    tbl = table_shape.table
    if headers:
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(16)
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    r.font.name = theme.get("header_font", "Microsoft YaHei")
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(theme.get("primary_color", "1E2761"))
    for ri, row in enumerate(rows):
        excel_row = ri + (1 if headers else 0)
        for ci, val in enumerate(row):
            cell = tbl.cell(excel_row, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(14)
                    r.font.color.rgb = hex_to_rgb("212121")
                    r.font.name = theme.get("body_font", "Microsoft YaHei")


def add_thank_you(prs: Presentation, theme: dict[str, str], s: dict[str, Any]) -> None:
    # type: ignore[reportGeneralTypeIssues]
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    fill = blank.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(theme.get("primary_color", "1E2761"))
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.33), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(
        p.add_run(),
        s["title"],
        72,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        font=theme.get("header_font", "Microsoft YaHei"),
    )
    if s.get("subtitle"):
        stx = blank.shapes.add_textbox(
            Inches(0.5), Inches(5.0), Inches(12.33), Inches(1)
        )
        stf = stx.text_frame
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        set_run(
            sp.add_run(),
            s["subtitle"],
            24,
            color=RGBColor(0xCA, 0xDC, 0xFC),
            font=theme.get("body_font", "Microsoft YaHei"),
        )


def add_timeline(prs: Presentation, theme: dict[str, str], s: dict[str, Any]) -> None:
    # type: ignore[reportGeneralTypeIssues]
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    tx = blank.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    set_run(
        p.add_run(),
        s["title"],
        32,
        bold=True,
        color=hex_to_rgb(theme.get("primary_color", "1E2761")),
        font=theme.get("header_font", "Microsoft YaHei"),
    )
    events = s.get("events", [])
    if not events:
        return
    n = len(events)
    step = 12.33 / max(n, 1)
    line = blank.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.6), Inches(12.33), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(theme.get("primary_color", "1E2761"))
    line.line.fill.background()
    for i, ev in enumerate(events):
        cx = Inches(0.5 + step * (i + 0.5))
        dot = blank.shapes.add_shape(
            MSO_SHAPE.OVAL, cx - Inches(0.15), Inches(3.45), Inches(0.3), Inches(0.3)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb(theme.get("accent_color", "F96167"))
        dot.line.fill.background()
        ltx = blank.shapes.add_textbox(
            cx - Inches(step / 2), Inches(2.0), Inches(step), Inches(1.3)
        )
        ltf = ltx.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        date_text = ev.get("date", "")
        title_text = ev.get("title", "")
        set_run(
            lp.add_run(),
            date_text,
            12,
            bold=True,
            color=hex_to_rgb(theme.get("accent_color", "F96167")),
            font=theme.get("body_font", "Microsoft YaHei"),
        )
        lp2 = ltf.add_paragraph()
        lp2.alignment = PP_ALIGN.CENTER
        set_run(
            lp2.add_run(),
            title_text,
            14,
            bold=True,
            color=hex_to_rgb(theme.get("primary_color", "1E2761")),
            font=theme.get("body_font", "Microsoft YaHei"),
        )
        if ev.get("description"):
            dtx = blank.shapes.add_textbox(
                cx - Inches(step / 2), Inches(4.0), Inches(step), Inches(2.0)
            )
            dtf = dtx.text_frame
            dtf.word_wrap = True
            dp = dtf.paragraphs[0]
            dp.alignment = PP_ALIGN.CENTER
            set_run(
                dp.add_run(),
                ev["description"],
                10,
                color=hex_to_rgb("212121"),
                font=theme.get("body_font", "Microsoft YaHei"),
            )


RENDERERS = {
    "title": add_title,
    "section": add_section,
    "bullet": add_bullet,
    "two_column": add_two_column,
    "three_column": add_two_column,
    "stat": add_stat,
    "quote": add_quote,
    "process": add_process,
    "table": add_table,
    "timeline": add_timeline,
    "thank_you": add_thank_you,
}


def main() -> None:
    parser = argparse.ArgumentParser(description="从 JSON spec 生成 PPTX")
    parser.add_argument("spec", help="spec.json 文件路径")
    parser.add_argument("output", help="输出 PPTX 文件路径")
    args = parser.parse_args()
    spec = load_json(Path(args.spec))
    out_path = Path(args.output)
    theme = spec.get("theme", {})

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    n = 0
    for slide in spec.get("slides", []):
        stype = slide.get("type", "bullet")
        renderer = RENDERERS.get(stype)
        if not renderer:
            logger.warning("未知类型 %s, 跳过", stype)
            continue
        renderer(prs, theme, slide)
        n += 1

    prs.save(str(out_path))
    logger.info("已生成: %s (%d 页)", out_path, n)


if __name__ == "__main__":
    main()
