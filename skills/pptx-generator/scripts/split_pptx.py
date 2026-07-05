#!/usr/bin/env python3
"""把一个 PPTX 按页拆成多个单页 PPTX
用法: split_pptx.py <input.pptx> <output_dir>
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import logging
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from skills.shared.logging_config import setup_logging

setup_logging()
logger = logging.getLogger(__name__)


def main() -> None:
    parser = argparse.ArgumentParser(description="拆分 PPTX 为单页文件")
    parser.add_argument("input", help="输入 PPTX 文件路径")
    parser.add_argument("output_dir", help="输出目录路径")
    args = parser.parse_args()

    src = Path(args.input)
    out_dir = Path(args.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(src))

    for idx, slide in enumerate(prs.slides, 1):
        new_prs = Presentation()
        new_prs.slide_width = prs.slide_width
        new_prs.slide_height = prs.slide_height
        blank_layout = new_prs.slide_layouts[6]
        new_slide = new_prs.slides.add_slide(blank_layout)
        for shape in slide.shapes:
            el = deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(el, "p:extLst")
        out_file = out_dir / f"slide_{idx:03d}.pptx"
        new_prs.save(str(out_file))

    logger.info("已拆分 %d 页到 %s", len(prs.slides), out_dir)


if __name__ == "__main__":
    main()
