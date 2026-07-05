#!/usr/bin/env python3
"""估算演讲时长(基于 speaker_notes.md 或 deck.pptx 的备注)

用法:
    estimate_duration.py notes.md
    estimate_duration.py deck.pptx
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import logging
import re
import sys
from pathlib import Path

from skills.shared.logging_config import setup_logging

setup_logging()
logger = logging.getLogger(__name__)


def estimate_from_md(md_path):
    text = Path(md_path).read_text(encoding="utf-8-sig")
    sections = re.split(r"^##\s+页\s+\d+:", text, flags=re.MULTILINE)
    sections = [s for s in sections if s.strip()]
    total = 0.0
    logger.info("%-4s %-8s %-10s", "页", "字数", "时长(分钟)")
    logger.info("-" * 30)
    for i, sec in enumerate(sections, 1):
        clean = re.sub(r"[*#>_-]+", "", sec)
        cn_chars = len(re.findall(r"[\u4e00-\u9fff]", clean))
        en_words = len(re.findall(r"\b[a-zA-Z]+\b", clean))
        equiv = cn_chars + en_words * 0.6
        minutes = equiv / 165
        total += minutes
        logger.info("%-4d %-8d %.2f", i, int(equiv), minutes)
    logger.info("-" * 30)
    logger.info("总时长: %.2f 分钟(不含 Q&A)", total)
    return total


def estimate_from_pptx(pptx_path):
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    total = 0.0
    logger.info("%-4s %-8s %-10s", "页", "字数", "时长(分钟)")
    logger.info("-" * 30)
    for i, slide in enumerate(prs.slides, 1):
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
        clean = re.sub(r"[*#>_-]+", "", notes)
        cn = len(re.findall(r"[\u4e00-\u9fff]", clean))
        en = len(re.findall(r"\b[a-zA-Z]+\b", clean))
        equiv = cn + en * 0.6
        minutes = equiv / 165
        total += minutes
        logger.info("%-4d %-8d %.2f", i, int(equiv), minutes)
    logger.info("-" * 30)
    logger.info("总时长: %.2f 分钟(不含 Q&A)", total)
    return total


def main() -> None:
    parser = argparse.ArgumentParser(description="估算演讲时长")
    parser.add_argument("input", help="输入文件路径 (notes.md 或 deck.pptx)")
    args = parser.parse_args()
    path = Path(args.input)
    if path.suffix == ".md":
        estimate_from_md(path)
    elif path.suffix == ".pptx":
        estimate_from_pptx(path)
    else:
        logger.error("不支持的文件类型: %s", path.suffix)
        sys.exit(1)


if __name__ == "__main__":
    main()
