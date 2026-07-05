#!/usr/bin/env python3
"""把 speaker_notes.md 写入 PPTX 的演讲者备注栏。

speaker_notes.md 格式:
    ## 页 1: <标题>
    **承接**: ...
    **展开**: ...
    **落地**: ...
    **转场到下一页**: ...
    **时长**: 1.5 分钟
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import logging
import re
import sys
from pathlib import Path

from pptx import Presentation
from skills.shared.logging_config import setup_logging

setup_logging()
logger = logging.getLogger(__name__)


SECTION_PATTERN = re.compile(r"^##\s+页\s+(\d+):\s*(.+)$", re.MULTILINE)


def parse_notes(md_text):
    sections = SECTION_PATTERN.split(md_text)
    notes = {}
    for i in range(1, len(sections), 3):
        page_num = int(sections[i])
        title = sections[i + 1].strip()
        body = sections[i + 2].strip() if i + 2 < len(sections) else ""
        notes[page_num] = {"title": title, "body": body}
    return notes


def main() -> None:
    parser = argparse.ArgumentParser(description="把演讲者备注写入 PPTX")
    parser.add_argument("deck", help="PPTX 文件路径")
    parser.add_argument("notes", help="备注 Markdown 文件路径")
    parser.add_argument("output", nargs="?", help="输出 PPTX 文件路径")
    args = parser.parse_args()
    deck_path = Path(args.deck)
    notes_path = Path(args.notes)
    out_path = Path(args.output) if args.output else deck_path

    notes = parse_notes(notes_path.read_text(encoding="utf-8-sig"))
    prs = Presentation(str(deck_path))

    n = 0
    for idx, slide in enumerate(prs.slides, 1):
        if idx in notes:
            slide.notes_slide.notes_text_frame.text = notes[idx]["body"]
            n += 1
            logger.info("页 %d: 已写入备注", idx)

    prs.save(str(out_path))
    logger.info("已写入 %d 页备注到: %s", n, out_path)


if __name__ == "__main__":
    main()
