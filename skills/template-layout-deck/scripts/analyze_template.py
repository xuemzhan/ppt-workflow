#!/usr/bin/env python3
"""分析 PPTX 模板结构,输出 JSON 配置
用法: analyze_template.py <template.pptx> [output.json]
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import json
import logging
import sys
from pathlib import Path

from pptx import Presentation
from skills.shared.logging_config import setup_logging

setup_logging()
logger = logging.getLogger(__name__)


def main() -> None:
    parser = argparse.ArgumentParser(description="分析 PPTX 模板结构")
    parser.add_argument("template", help="模板 PPTX 文件路径")
    parser.add_argument("output", nargs="?", help="输出 JSON 文件路径")
    args = parser.parse_args()

    tmpl = Presentation(args.template)
    info = {
        "template_path": args.template,
        "slide_size": {"width": tmpl.slide_width, "height": tmpl.slide_height},
        "layouts": [],
    }
    for i, layout in enumerate(tmpl.slide_layouts):
        phs = []
        for ph in layout.placeholders:
            phs.append(
                {
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                }
            )
        info["layouts"].append({"index": i, "name": layout.name, "placeholders": phs})

    out = json.dumps(info, ensure_ascii=False, indent=2)
    if args.output:
        Path(args.output).write_text(out, encoding="utf-8")
        logger.info("已写入: %s", args.output)
    else:
        print(out)


if __name__ == "__main__":
    main()
