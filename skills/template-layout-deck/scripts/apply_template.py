#!/usr/bin/env python3
"""模板套用:把内容 spec 应用到指定 .pptx 模板的版式上

用法:
    python apply_template.py <template.pptx> <content.json> <output.pptx>
    python apply_template.py <template.pptx> --analyze

content.json 格式:
{
  "slides": [
    {
      "layout_index": 0,
      "layout_name": "封面",
      "placeholders": {"0": "标题文本", "1": "副标题文本"}
    }
  ]
}
"""

# IMPORTANT: path setup must come BEFORE any skills.* imports
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent.parent))

import argparse
import json
import logging

from skills.shared.logging_config import setup_logging
from skills.shared.utils import load_json

setup_logging()
logger = logging.getLogger(__name__)

from pptx import Presentation


def analyze_template(tmpl_path):
    tmpl = Presentation(str(tmpl_path))
    info = {
        "slide_size": {"width": tmpl.slide_width, "height": tmpl.slide_height},
        "layouts": [],
    }
    for i, layout in enumerate(tmpl.slide_layouts):
        phs = []
        for ph in layout.placeholders:
            phs.append(
                {
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                }
            )
        info["layouts"].append({"index": i, "name": layout.name, "placeholders": phs})
    return info


def add_with_layout(deck, layout_index, content):
    layout = deck.slide_layouts[layout_index]
    slide = deck.slides.add_slide(layout)
    for ph in slide.placeholders:
        idx = str(ph.placeholder_format.idx)
        if idx in content:
            ph.text = content[idx]
    return slide


def main() -> None:
    parser = argparse.ArgumentParser(
        description="模板套用: 把内容 spec 应用到 .pptx 模板"
    )
    parser.add_argument("template", help="模板 .pptx 文件路径")
    parser.add_argument("--analyze", action="store_true", help="分析模板版式")
    parser.add_argument("content", nargs="?", help="内容 JSON 文件路径")
    parser.add_argument("output", nargs="?", help="输出 .pptx 文件路径")
    args = parser.parse_args()
    tmpl_path = Path(args.template)

    if args.analyze:
        info = analyze_template(tmpl_path)
        print(json.dumps(info, ensure_ascii=False, indent=2))
        return

    if not args.content or not args.output:
        parser.error("缺少 content.json 或 output.pptx")

    content_path = Path(args.content)
    out_path = Path(args.output)
    content = load_json(content_path)

    deck = Presentation(str(tmpl_path))

    n = 0
    for slide_spec in content.get("slides", []):
        layout_idx = slide_spec.get("layout_index", 1)
        if layout_idx >= len(deck.slide_layouts):
            logger.warning("版式索引 %d 超出范围, 跳过", layout_idx)
            continue
        add_with_layout(deck, layout_idx, slide_spec.get("placeholders", {}))
        n += 1

    deck.save(str(out_path))
    logger.info("已套用模板并生成: %s (%d 页)", out_path, n)


if __name__ == "__main__":
    main()
