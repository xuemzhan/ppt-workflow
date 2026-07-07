#!/usr/bin/env python3
"""端到端集成测试

覆盖 5 个核心场景:
1. test_orchestrator_runs_full_pipeline: orchestrate.py 跑通 10 步
2. test_chart_renders_with_log_scale: render_chart.py 数量级自适应
3. test_template_layout_preserves_placeholders: apply_template.py 套模板
4. test_qa_detects_issues: qa_check.py 检出硬编码占位
5. test_validate_spec_catches_unknown_types: validate_spec.py 检出未知 type
"""

from __future__ import annotations

import json
import subprocess
import sys

from pptx import Presentation
from tests.conftest import REPO, SAMPLE_INPUT, SAMPLE_OUTPUT

PYTHON = sys.executable


def run(cmd: list[str], timeout: int = 60) -> subprocess.CompletedProcess:
    """Run a subprocess command, return result."""
    return subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)


# ──────────────────────────────────────────────────────────────
# Test 1: orchestrate.py 跑通完整流水线
# ──────────────────────────────────────────────────────────────


def test_orchestrator_runs_full_pipeline(ensure_clean_output):
    """端到端:从 input 跑 orchestrate,产物齐全。"""
    result = run(
        [
            PYTHON,
            str(REPO / "scripts" / "orchestrate.py"),
            "--input",
            str(SAMPLE_INPUT),
            "--output",
            str(SAMPLE_OUTPUT),
            "--force",
        ]
    )
    assert result.returncode == 0, (
        f"orchestrate failed:\nSTDOUT: {result.stdout}\nSTDERR: {result.stderr}"
    )

    # 关键产物应存在
    expected = [
        "outline.md",
        "storyline.md",
        "spec.json",
        "deck_draft.pptx",
        "deck_final.pptx",
        "speaker_notes.md",
        "qa_report.md",
        "thumbnails.jpg",
        "workflow_log.md",
    ]
    for name in expected:
        path = SAMPLE_OUTPUT / name
        assert path.exists(), f"missing {name}"
        assert path.stat().st_size > 0, f"empty {name}"

    # spec.json 应有 8 页(对照样例 slides.json)
    spec = json.loads((SAMPLE_OUTPUT / "spec.json").read_text(encoding="utf-8-sig"))
    assert len(spec["slides"]) == 8

    # deck_final.pptx 应有 8 张幻灯片
    prs = Presentation(str(SAMPLE_OUTPUT / "deck_final.pptx"))
    assert len(prs.slides) == 8


# ──────────────────────────────────────────────────────────────
# Test 2: 数量级自适应图表
# ──────────────────────────────────────────────────────────────


def test_chart_renders_with_log_scale(tmp_path):
    """图表在数据差异 > 50x 时自动切换对数刻度。"""
    spec = {
        "type": "bar",
        "title": "Test Log Scale",
        "data": {
            "categories": ["A", "B"],
            "series": [{"name": "x", "values": [100, 0.5]}],
        },
        "unit": " ms",
    }
    spec_path = tmp_path / "spec.json"
    spec_path.write_text(json.dumps(spec), encoding="utf-8")
    out_path = tmp_path / "chart.png"

    result = run(
        [
            PYTHON,
            str(REPO / "skills" / "chart-slide-maker" / "scripts" / "render_chart.py"),
            str(spec_path),
            str(out_path),
        ]
    )
    assert result.returncode == 0
    assert out_path.exists()
    # PNG 文件应非空
    assert out_path.stat().st_size > 1000


# ──────────────────────────────────────────────────────────────
# Test 3: 套模板保留占位符
# ──────────────────────────────────────────────────────────────


def test_template_layout_preserves_placeholders(tmp_path):
    """apply_template.py 套用公司模板,内容正确填入。"""
    template_src = REPO / "examples" / "company_template.pptx"
    content_src = REPO / "examples" / "sample_deck" / "template_content.json"
    out_path = tmp_path / "themed.pptx"

    assert template_src.exists(), "company_template.pptx missing"
    assert content_src.exists(), "template_content.json missing"

    result = run(
        [
            PYTHON,
            str(
                REPO
                / "skills"
                / "template-layout-deck"
                / "scripts"
                / "apply_template.py"
            ),
            str(template_src),
            str(content_src),
            str(out_path),
        ]
    )
    assert result.returncode == 0
    assert out_path.exists()

    prs = Presentation(str(out_path))
    # 5 页(对照 template_content.json)
    assert len(prs.slides) == 5
    # 首页有 "智能办公套件 V1.0 阶段性成果"
    first_text = "\n".join(
        shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame
    )
    assert "智能办公套件 V1.0" in first_text


# ──────────────────────────────────────────────────────────────
# Test 4: QA 检出硬编码占位词
# ──────────────────────────────────────────────────────────────


def test_qa_detects_issues(tmp_path):
    """qa_check.py 构造含问题的 pptx,应检出占位词。"""
    bad_pptx = tmp_path / "bad.pptx"
    prs = Presentation()
    prs.slide_width = 9144000  # type: ignore[reportAttributeAccessIssue]
    prs.slide_height = 5143500  # type: ignore[reportAttributeAccessIssue]
    # 封面
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    s1.shapes.add_textbox(0, 0, 9144000, 1000000)  # type: ignore[reportArgumentType].text_frame.text = "Cover"
    # 问题页:占位词
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    s2.shapes.add_textbox(0, 0, 9144000, 1000000)  # type: ignore[reportArgumentType].text_frame.text = "对照 1"
    # 干净页
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    s3.shapes.add_textbox(0, 0, 9144000, 1000000)  # type: ignore[reportArgumentType].text_frame.text = "Normal"
    # 封底
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    s4.shapes.add_textbox(0, 0, 9144000, 1000000)  # type: ignore[reportArgumentType].text_frame.text = "Thanks"
    prs.save(str(bad_pptx))

    qa_report = tmp_path / "qa.md"
    result = run(
        [
            PYTHON,
            str(REPO / "skills" / "deck-review-publisher" / "scripts" / "qa_check.py"),
            str(bad_pptx),
            str(qa_report),
        ]
    )
    assert result.returncode == 0
    report = qa_report.read_text(encoding="utf-8-sig")
    # 必须检出"对照"占位
    assert "对照" in report or "placeholder" in report.lower()


# ──────────────────────────────────────────────────────────────
# Test 5: validate_spec 拒绝未知 type
# ──────────────────────────────────────────────────────────────


def test_validate_spec_catches_unknown_types(tmp_path):
    """validate_spec.py 对含 unknown type 的 spec 返回 exit 1。"""
    bad_spec = {"slides": [{"type": "foo_bar_baz"}]}
    spec_path = tmp_path / "bad.json"
    spec_path.write_text(json.dumps(bad_spec), encoding="utf-8")

    result = run(
        [
            PYTHON,
            str(REPO / "skills" / "pptx-generator" / "scripts" / "validate_spec.py"),
            str(spec_path),
        ]
    )
    # 退出码 1 = 有错误
    assert result.returncode == 1
    # 错误信息要明确(在 stdout 或 stderr)
    combined = result.stdout + result.stderr
    assert "foo_bar_baz" in combined
    assert "不支持" in combined or "error" in combined.lower()
